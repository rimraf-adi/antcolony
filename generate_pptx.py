"""Generate a clean black-and-white PowerPoint presentation for the ACO Edge Detection mini project."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ──
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(100, 100, 100)
LIGHT_GRAY = RGBColor(200, 200, 200)
DARK_GRAY = RGBColor(50, 50, 50)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=16, bold=False, color=BLACK,
                  alignment=PP_ALIGN.LEFT, font_name="Arial",
                  space_before=Pt(4), space_after=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(1.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.fill.background()


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        add_textbox(slide, left, top, Inches(3), Inches(0.4),
                    f"[Image not found: {os.path.basename(path)}]",
                    font_size=10, color=GRAY)


def bullet_list(tf, items, font_size=14, color=BLACK, indent_level=0):
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.level = indent_level
        p.space_before = Pt(4)
        p.space_after = Pt(2)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]  # blank slide

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                "ACO-Based Edge Detection", font_size=42, bold=True,
                color=BLACK, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                "Ant Colony Optimization for Image Edge Detection\nEvaluated on the Berkeley Segmentation Dataset (BSDS300)",
                font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

    add_line(slide, Inches(4), Inches(3.8), Inches(5))

    info_tf = add_textbox(slide, Inches(2.5), Inches(4.2), Inches(8), Inches(2.5),
                          "Mini Project — Image Analysis & Computer Vision",
                          font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    add_paragraph(info_tf, "", font_size=8)
    add_paragraph(info_tf, "Team Members: [Update Names Here]",
                  font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_paragraph(info_tf, "Guide: [Professor Name]",
                  font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_paragraph(info_tf, "Institution: [College / University Name]",
                  font_size=14, color=BLACK, alignment=PP_ALIGN.CENTER)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 2: Problem Statement
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Problem Statement", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    tf = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                     "Edge detection is a fundamental operation in image analysis — identifying "
                     "boundaries where pixel intensity changes sharply.",
                     font_size=16, color=BLACK)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "Challenges with Traditional Methods:", font_size=18,
                  bold=True, color=BLACK)
    bullet_list(tf, [
        "• Fixed-kernel methods (Sobel, Canny) cannot adaptively respond to varying edge strengths",
        "• Pixel-wise operators produce fragmented, disconnected edge contours",
        "• Sensitive to noise without careful parameter tuning",
    ], font_size=14)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "Our Approach:", font_size=18, bold=True, color=BLACK)
    bullet_list(tf, [
        "• Apply Ant Colony Optimization (ACO) — a swarm intelligence metaheuristic",
        "• Artificial ants traverse the pixel grid, depositing pheromone on edge locations",
        "• Pheromone accumulates along true edge contours over multiple iterations",
        "• Evaluate against human-annotated ground truth from BSDS300 (300 images)",
    ], font_size=14)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 3: Literature Survey
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Literature Survey", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    tf = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                     "[1]  Tian, Yu & Xie (2008) — IEEE Congress on Evolutionary Computation",
                     font_size=16, bold=True, color=BLACK)
    bullet_list(tf, [
        "    Pioneered ACO for image edge detection. Ants traverse pixel grids using",
        "    pheromone and gradient-based heuristics. Competitive with Sobel and Canny.",
    ], font_size=13, color=GRAY)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "[2]  Nezamabadi-pour et al. — Ant Colonies for MRF-Based Image Segmentation",
                  font_size=16, bold=True, color=BLACK)
    bullet_list(tf, [
        "    Extended ACO to MRF-based segmentation tasks, showing pheromone-guided",
        "    traversal captures spatial context beyond local gradient information.",
    ], font_size=13, color=GRAY)

    add_paragraph(tf, "", font_size=8)
    add_paragraph(tf, "[3]  IJETT V71I9P233 — Edge Detection Using Ant Colony Optimization",
                  font_size=16, bold=True, color=BLACK)
    bullet_list(tf, [
        "    Comprehensive survey comparing ACO with classical methods. Analyzes",
        "    parameter sensitivity (α, β, ρ). ACO produces smoother, more connected edges.",
    ], font_size=13, color=GRAY)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 4: Biological Inspiration
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Biological Inspiration", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    # Left column - Real Ants
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                "Real Ant Behavior", font_size=20, bold=True, color=BLACK)
    tf_l = add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.5),
                       "", font_size=14)
    bullet_list(tf_l, [
        "• Ants deposit pheromone on paths while foraging",
        "• Shorter, better paths accumulate more pheromone",
        "• Other ants prefer pheromone-rich paths",
        "• Old pheromone evaporates over time",
        "• Collective behavior → emergent optimization",
    ], font_size=14)

    # Right column - Our Analogy
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                "Our Analogy", font_size=20, bold=True, color=BLACK)
    tf_r = add_textbox(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(4.5),
                       "", font_size=14)
    bullet_list(tf_r, [
        "• Image pixels = graph nodes",
        "• 8-connected neighbors = possible moves",
        "• Gradient magnitude = heuristic desirability",
        "• Pheromone matrix = collective edge memory",
        "• Edge contours emerge from pheromone accumulation",
    ], font_size=14)

    # Divider line
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(6.5), Inches(1.5), Pt(1), Inches(4.5))
    div.fill.solid()
    div.fill.fore_color.rgb = LIGHT_GRAY
    div.line.fill.background()

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 5: Mathematical Formulation
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Mathematical Formulation", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
                "Transition Probability", font_size=20, bold=True, color=BLACK)

    add_textbox(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.5),
                "Each ant at pixel (i,j) chooses its next position from the 8-connected neighborhood:",
                font_size=14, color=GRAY)

    # Formula box
    formula_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2.7), Inches(9), Inches(0.8)
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    formula_box.line.color.rgb = LIGHT_GRAY
    formula_box.text_frame.paragraphs[0].text = "P(i,j)  =  [ τ(i,j)^α  ·  η(i,j)^β ]  /  Σ  [ τ(m,n)^α  ·  η(m,n)^β ]"
    formula_box.text_frame.paragraphs[0].font.size = Pt(16)
    formula_box.text_frame.paragraphs[0].font.name = "Courier New"
    formula_box.text_frame.paragraphs[0].font.color.rgb = BLACK
    formula_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5),
                "Pheromone Update", font_size=20, bold=True, color=BLACK)

    formula_box2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(4.4), Inches(9), Inches(0.8)
    )
    formula_box2.fill.solid()
    formula_box2.fill.fore_color.rgb = RGBColor(245, 245, 245)
    formula_box2.line.color.rgb = LIGHT_GRAY
    formula_box2.text_frame.paragraphs[0].text = "τ(i,j)  ←  (1 − ρ) · τ(i,j)  +  Δτ(i,j)"
    formula_box2.text_frame.paragraphs[0].font.size = Pt(16)
    formula_box2.text_frame.paragraphs[0].font.name = "Courier New"
    formula_box2.text_frame.paragraphs[0].font.color.rgb = BLACK
    formula_box2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Symbol definitions
    tf_sym = add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.8),
                         "Where:", font_size=14, bold=True, color=BLACK)
    bullet_list(tf_sym, [
        "τ(i,j) = pheromone at pixel (i,j)     η(i,j) = heuristic value (Sobel gradient magnitude)",
        "α = pheromone influence weight          β = heuristic influence weight          ρ = evaporation rate",
    ], font_size=13, color=GRAY)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 6: System Architecture
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "System Architecture & Pipeline", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    # Flowchart steps
    steps = [
        ("1", "Input Image\n(Grayscale)"),
        ("2", "Sobel Gradient\nη(i,j)"),
        ("3", "Init Pheromone\nτ(i,j) = 0.1"),
        ("4", "Ant Traversal\nK ants, S steps"),
        ("5", "Pheromone\nUpdate"),
        ("6", "Edge Map\n[0, 255]"),
    ]

    box_w = Inches(1.7)
    box_h = Inches(1.1)
    start_x = Inches(0.6)
    y_pos = Inches(1.8)
    gap = Inches(0.35)

    for i, (num, label) in enumerate(steps):
        x = start_x + i * (box_w + gap)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BLACK
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Step {num}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(10)
        p2.font.color.rgb = GRAY
        p2.font.name = "Arial"
        p2.alignment = PP_ALIGN.CENTER

        # Arrow
        if i < len(steps) - 1:
            arrow_x = x + box_w
            add_textbox(slide, arrow_x, y_pos + Inches(0.3), gap, Inches(0.4),
                        "→", font_size=20, color=BLACK, alignment=PP_ALIGN.CENTER)

    # Software stack
    add_textbox(slide, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.5),
                "Software Stack", font_size=20, bold=True, color=BLACK)

    tf_sw = add_textbox(slide, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3),
                        "", font_size=14)
    bullet_list(tf_sw, [
        "• Python ≥ 3.13",
        "• NumPy ≥ 2.4 — array operations",
        "• OpenCV ≥ 4.13 — image I/O, Sobel, Canny",
        "• tqdm — progress bars",
        "• uv — package manager",
    ], font_size=13)

    add_textbox(slide, Inches(7), Inches(4.0), Inches(5.5), Inches(3),
                "Code Structure", font_size=16, bold=True, color=BLACK)
    tf_code = add_textbox(slide, Inches(7), Inches(4.5), Inches(5.5), Inches(2.5),
                          "", font_size=13)
    bullet_list(tf_code, [
        "• main.py — AntColonyEdgeDetector class",
        "• process_images.py — Batch pipeline + CLI",
    ], font_size=13)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 7: Implementation (Code)
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Core Algorithm — Code", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    code_text = """class AntColonyEdgeDetector:
    def run_detection(self):
        self._initialize_matrices()       # Sobel gradient + pheromone init

        for iteration in range(self.num_iterations):
            # Place K ants at random pixel positions
            ant_rows = np.random.randint(0, H, self.num_ants)
            ant_cols = np.random.randint(0, W, self.num_ants)

            for step in range(steps_per_ant):
                for ant in range(self.num_ants):
                    # Transition probability: P = τ^α · η^β
                    # Choose next pixel from 8-connected neighbors
                    # Deposit pheromone proportional to gradient

            # Global pheromone update with evaporation
            τ = (1 - ρ) * τ + Δτ

        # Normalize pheromone matrix → edge map [0, 255]
        # Dilate edges with 3×3 kernel for visibility
        return dilated_edge_map"""

    code_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(11), Inches(5.3)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    code_box.line.color.rgb = LIGHT_GRAY
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(12)
    p.font.name = "Courier New"
    p.font.color.rgb = BLACK
    p.line_spacing = Pt(17)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 8: Dataset
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Dataset — BSDS300", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
                "Berkeley Segmentation Dataset — 300 natural images with human-annotated ground truth",
                font_size=16, color=GRAY)

    # Stats boxes
    stats = [("300", "Natural Images"), ("200", "Training Set"),
             ("100", "Test Set"), ("5–7", "Annotators / Image")]
    for i, (val, label) in enumerate(stats):
        x = Inches(1) + i * Inches(2.8)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, Inches(2.3), Inches(2.3), Inches(1.2))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BLACK
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLACK
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(12)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    # Features
    add_textbox(slide, Inches(0.8), Inches(3.9), Inches(5.5), Inches(0.5),
                "Dataset Features", font_size=18, bold=True, color=BLACK)
    tf_l = add_textbox(slide, Inches(0.8), Inches(4.4), Inches(5.5), Inches(3),
                       "", font_size=13)
    bullet_list(tf_l, [
        "• Diverse scenes: animals, landscapes, people, objects",
        "• Multiple independent human segmentations per image",
        "• Standard benchmark in edge detection literature",
    ], font_size=13)

    add_textbox(slide, Inches(7), Inches(3.9), Inches(5.5), Inches(0.5),
                "Our Processing", font_size=18, bold=True, color=BLACK)
    tf_r = add_textbox(slide, Inches(7), Inches(4.4), Inches(5.5), Inches(3),
                       "", font_size=13)
    bullet_list(tf_r, [
        "• All images converted to grayscale",
        "• .seg files parsed → binary edge maps via Sobel",
        "• 3-panel dashboard visualization generated per image",
    ], font_size=13)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 9: Hyperparameters
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Hyperparameters", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    # Table
    rows, cols = 7, 5
    table_shape = slide.shapes.add_table(rows, cols,
                                          Inches(1), Inches(1.6),
                                          Inches(11), Inches(3.5))
    table = table_shape.table

    headers = ["Parameter", "Symbol", "Default", "Batch", "Description"]
    data = [
        ["num_ants", "K", "2,048", "8,192", "Ants per iteration"],
        ["num_iterations", "T", "20", "50", "ACO iterations"],
        ["alpha", "α", "1.0", "1.0", "Pheromone influence weight"],
        ["beta", "β", "2.0", "2.0", "Heuristic (gradient) influence weight"],
        ["decay_rate", "ρ", "0.1", "0.1", "Pheromone evaporation rate"],
        ["steps_per_ant", "S", "10", "10", "Steps each ant takes per iteration"],
    ]

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Arial"
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK

    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = BLACK
                p.font.name = "Courier New" if j == 0 else "Arial"

    add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(1),
                "Key insight: β > α favors gradient-guided exploration, producing edges "
                "aligned with true intensity transitions. Batch processing uses 8192 ants "
                "and 50 iterations for production-quality edge maps.",
                font_size=13, color=GRAY)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 10: Results 1
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Results — Sample Outputs (Set 1)", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4),
                "Each dashboard:  Original  →  Human Ground Truth  →  ACO Edge Map",
                font_size=14, color=GRAY)

    img1 = os.path.join(SCRIPT_DIR, "presentation", "assets", "sample_bears.png")
    add_image_safe(slide, img1, Inches(0.8), Inches(1.9), width=Inches(11.5))

    add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.3),
                "BSDS300 #100075 — Bears at water's edge", font_size=11, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    img2 = os.path.join(SCRIPT_DIR, "presentation", "assets", "sample_plane.png")
    add_image_safe(slide, img2, Inches(0.8), Inches(4.8), width=Inches(11.5))

    add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.3),
                "BSDS300 #3096 — Military aircraft in flight", font_size=11, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 11: Results 2
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Results — Sample Outputs (Set 2)", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    img3 = os.path.join(SCRIPT_DIR, "presentation", "assets", "sample_3.png")
    add_image_safe(slide, img3, Inches(0.8), Inches(1.5), width=Inches(11.5))

    add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.3),
                "BSDS300 #42049 — Bird on branch", font_size=11, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    img4 = os.path.join(SCRIPT_DIR, "presentation", "assets", "sample_4.png")
    add_image_safe(slide, img4, Inches(0.8), Inches(4.5), width=Inches(11.5))

    add_textbox(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.3),
                "BSDS300 #135069 — Eagles in flight", font_size=11, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 12: Analysis
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Strengths & Limitations", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    # Strengths
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                "Strengths", font_size=20, bold=True, color=BLACK)
    tf_s = add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(5),
                       "", font_size=13)
    bullet_list(tf_s, [
        "• Adaptive: pheromone reinforcement concentrates on true edges",
        "• Connected edges: ant traversal produces continuous contours",
        "• Noise tolerance: collective behavior averages out noise",
        "• Flexible: tunable via α, β, ρ, ant count parameters",
        "• No training required: pure optimization approach",
    ], font_size=13)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(6.5), Inches(1.5), Pt(1), Inches(5))
    div.fill.solid()
    div.fill.fore_color.rgb = LIGHT_GRAY
    div.line.fill.background()

    # Limitations
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                "Limitations", font_size=20, bold=True, color=BLACK)
    tf_w = add_textbox(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(5),
                       "", font_size=13)
    bullet_list(tf_w, [
        "• Computational cost: per-ant, per-step loops are expensive",
        "• Parameter sensitivity: poor α/β ratios produce weak edges",
        "• Thinner edges: ACO edges sparser than ground truth",
        "• No multi-scale: single-resolution Sobel heuristic",
        "• Sequential bottleneck: hard to parallelize on CPU",
    ], font_size=13)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 13: Comparison Table
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "ACO vs. Traditional Methods", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    rows2, cols2 = 7, 4
    table_shape2 = slide.shapes.add_table(rows2, cols2,
                                           Inches(1), Inches(1.6),
                                           Inches(11), Inches(4))
    table2 = table_shape2.table

    headers2 = ["Property", "Sobel", "Canny", "ACO (Ours)"]
    data2 = [
        ["Approach", "Fixed 3×3 kernels", "Multi-stage pipeline", "Swarm intelligence"],
        ["Adaptivity", "None", "Threshold-based", "Pheromone-guided"],
        ["Edge Continuity", "Low (pixel-wise)", "Medium (hysteresis)", "High (ant paths)"],
        ["Noise Handling", "Poor", "Good (Gaussian blur)", "Good (collective avg.)"],
        ["Speed", "Very fast", "Fast", "Slow (iterative)"],
        ["Parameters", "Kernel size", "2 thresholds + σ", "α, β, ρ, K, T, S"],
    ]

    for j, h in enumerate(headers2):
        cell = table2.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Arial"
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK

    for i, row_data in enumerate(data2):
        for j, val in enumerate(row_data):
            cell = table2.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = BLACK
                p.font.name = "Arial"

    add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(1),
                "Key takeaway: ACO trades computational speed for adaptivity and edge continuity. "
                "Best suited for applications where edge quality matters more than real-time performance.",
                font_size=13, color=GRAY)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 14: Conclusion & Future Work
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "Conclusion & Future Work", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    # Conclusions
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
                "Conclusions", font_size=20, bold=True, color=BLACK)
    tf_c = add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(5),
                       "", font_size=13)
    bullet_list(tf_c, [
        "• Successfully implemented ACO edge detection",
        "• Evaluated on all 300 BSDS300 images",
        "• Produces coherent, continuous edge contours",
        "• Fully unsupervised — no training or labels required",
        "• Dashboard visualizations enable direct comparison",
        "  with human ground truth",
    ], font_size=13)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(6.5), Inches(1.5), Pt(1), Inches(5))
    div.fill.solid()
    div.fill.fore_color.rgb = LIGHT_GRAY
    div.line.fill.background()

    # Future Work
    add_textbox(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.5),
                "Future Work", font_size=20, bold=True, color=BLACK)
    tf_f = add_textbox(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(5),
                       "", font_size=13)
    bullet_list(tf_f, [
        "• GPU acceleration via CUDA/OpenCL",
        "• Multi-scale heuristics using Gaussian pyramids",
        "• Hybrid ACO + Canny for thin, precise edges",
        "• Quantitative metrics: F1, ODS, OIS scores",
        "• Adaptive parameter tuning via PSO or",
        "  genetic algorithms",
    ], font_size=13)

    # ═════════════════════════════════════════════════════════════════
    # SLIDE 15: References + Thank You
    # ═════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
                "References", font_size=32, bold=True, color=BLACK)
    add_line(slide, Inches(0.8), Inches(1.15), Inches(4))

    refs = [
        "[1]  Tian, J., Yu, W., & Xie, S. (2008). \"An Ant Colony Optimization Algorithm for Image Edge Detection.\" IEEE Congress on Evolutionary Computation, pp. 751–756.",
        "[2]  Nezamabadi-pour, H., et al. \"Ant Colonies for MRF-Based Image Segmentation.\" Journal of Applied Sciences.",
        "[3]  \"Edge Detection Using Ant Colony Optimization.\" International Journal of Engineering Trends and Technology (IJETT), V71I9P233.",
        "[4]  Martin, D., Fowlkes, C., Tal, D., & Malik, J. (2001). \"A Database of Human Segmented Natural Images.\" ICCV.",
        "[5]  Dorigo, M. & Stützle, T. (2004). Ant Colony Optimization. MIT Press.",
        "[6]  Canny, J. (1986). \"A Computational Approach to Edge Detection.\" IEEE TPAMI, 8(6), pp. 679–698.",
    ]

    tf_ref = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(4),
                         refs[0], font_size=12, color=BLACK)
    for r in refs[1:]:
        add_paragraph(tf_ref, r, font_size=12, color=BLACK,
                      space_before=Pt(8), space_after=Pt(4))

    # Thank you
    add_textbox(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.7),
                "Thank You", font_size=28, bold=True, color=BLACK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                "Questions & Discussion", font_size=16, color=GRAY,
                alignment=PP_ALIGN.CENTER)

    # ── Save ──
    output_path = os.path.join(SCRIPT_DIR, "ACO_Edge_Detection_Presentation.pptx")
    prs.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")
    print(f"  {len(prs.slides)} slides, landscape 13.33\" × 7.5\"")


if __name__ == "__main__":
    create_presentation()
